from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat,RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import random
from os.path import dirname, realpath
from pptx.enum.lang import MSO_LANGUAGE_ID
import os
import requests
from get_image import get_image,translator
filepath = realpath(__file__)
design=["template1.pptx","template1.pptx"]
from extract.content import get_content




class MyImage:
    def __init__(self,img_name,prs,lang):
       
       #aux_im = Image.open(io.BytesIO(r.content))
       blank_slide_layout = prs.slide_layouts[0]
       slide = prs.slides.add_slide(blank_slide_layout)
       left = Inches(2)
       top=Inches(2.5)
       
       pic = slide.shapes.add_picture(img_name,left=left,top=top,width=Inches(7), height=Inches(3.5))
       line = pic.line
       line.color.rgb = RGBColor(0,0,0)



class MyMain:
    def __init__(self,data,prs,lang):
        self.layout=prs.slide_layouts[1]
        self.slide=prs.slides.add_slide(self.layout)
        self.title=self.slide.shapes.title
        if lang=='en':
            self.title.text=data['topic']
        else:
            self.title.text=translator(data['topic'],lang)["translated_text"]
        for shape in self.slide.shapes:
             
             if not shape.has_text_frame:
                 continue
             text_frame0 = self.slide.shapes[0].text_frame
             text_frame = self.slide.shapes[1].text_frame
             p = text_frame.add_paragraph()
             
             run=p.add_run()
             
             if lang=='en':
                run.text=data["introduction"]
             else:
                 run.text=translator(data["introduction"],lang)["translated_text"]
             font=run.font
             
             font.size=Pt(24)
             font.bold=True
             font.underline=True
             run.font.color.rgb  = RGBColor(0,100, 255)
             
             #font.language_id = MSO_LANGUAGE_ID.MIXED
             #font.color.theme_color= MSO_THEME_COLOR_INDEX.ACCENT_1
             text_frame.margin_left = 0
             text_frame.vertical_anchor = MSO_ANCHOR.TOP
             text_frame.word_wrap = False
             text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
             break


class MyContent:
     def __init__(self,slides,prs,lang):
         self.layout=prs.slide_layouts[2]
         self.slide=prs.slides.add_slide(self.layout)
         self.title=self.slide.shapes.title
         
         if lang=='en':
                self.title.text="Content"
         else:
               self.title.text=translator("Content",lang)["translated_text"]
         i=0
         for shape in self.slide.shapes:
             
             if not shape.has_text_frame:
                 continue
             text_frame = self.slide.shapes[1].text_frame
             i=1
             for l in slides['sections']:
                if not len(l['content'])==0:
                  x=str(i)+". "+l["heading"]
                  i=i+1
                  p = text_frame.add_paragraph()
                  run=p.add_run()
                  if lang=='en':
                     run.text=x
                  else:
                     run.text=translator(x,lang)["translated_text"]
                  font=run.font
                  font.size=Pt(18)
                  font.bold=True
                  #font.language_id = MSO_LANGUAGE_ID.MIXED
                  run.font.color.rgb  = RGBColor(0,0,139)
             break;


class MySlides:
     def __init__(self,key,value,prs,lang):
         self.layout=prs.slide_layouts[1]
         self.slide=prs.slides.add_slide(self.layout)
         self.title=self.slide.shapes.title
         self.title.text=key
         i=0
         for shape in self.slide.shapes:
             
             if not shape.has_text_frame:
                 continue
             text_frame = self.slide.shapes[1].text_frame
             index=0
             for l in value:
                  index+=1
                  x=l["title"]
                  y=l["desc"]
                  z=""
                  if y:
                      z=y[0]
                  if len(x)==0 and len(y)==0:
                    continue
                  
                  p = text_frame.add_paragraph()
                  run=p.add_run()
                  
                  if lang=='en':
                      if x and z :
                          x=str(index)+"."+x+": "
                          run.text=x
                      else:
                         x=str(index)+"."
                         run.text=x
                  else:
                    if x and z :
                          x=str(index)+"."+x+": "
                          run.text=translator(x,lang)["translated_text"]
                    else:
                         x=str(index)+"."
                         run.text=translator(x,lang)["translated_text"]

                  font=run.font
                  font.size=Pt(24)
                  font.name='Abyssinica SIL'
                  font.bold=True
                  #font.underline=True
                  font.color.rgb  = RGBColor(0,51,102)
                  #font.language_id = MSO_LANGUAGE_ID.MIXED
                  run = p.add_run()

                  font=run.font
                  font.size = Pt(18)
                  font.name='Calibri'
                  #font.language_id = MSO_LANGUAGE_ID.MIXED
                  try:

                    if lang=='en':
                       run.text =z
                    
                    else:
                      run.text =(translator(z,lang)["translated_text"])
                  except:
                    pass
                  font.color.rgb  = RGBColor(0,102,204)
                  #font.color.theme_color= MSO_THEME_COLOR_INDEX.ACCENT_1
                  
             break;



class MyExercise:
    def __init__(self,prs,lang):
         self.layout=prs.slide_layouts[2]
         self.slide=prs.slides.add_slide(self.layout)
         self.title=self.slide.shapes.title
         if lang=='en':
           self.title.text="Exercise"
         else:
            self.title.text=translator("Practise",lang)["translated_text"]
       
def generate_ppt(slides,lang='en'):
   prs = Presentation(random.choice(design))
   filename = slides['topic']+"-"+lang + '.pptx'
   loc = "files/ppts/"+filename
   img_count=0   
   MyMain(slides,prs,lang)
   MyContent(slides,prs,lang)
   images=[]    
   for d in slides["sections"]: 
     i=0
     while i<len(d['content']):
             
             MySlides(d['heading'],d['content'][i:i+4],prs,lang)

             i=i+4
     if d['img']:
         try:
             
             img=get_image(d['img'])
             MyImage(img,prs,lang)
             images.append(img)
             os.remove(d['img'])
         except:
           pass   
            
              
   MyExercise(prs,lang)
   prs.save(loc)
   for img in images:
     os.remove(img)
   print("ppt prepared!")
   print(slides)

       
      
   return loc

def get_ppt(content):




  if content is None:
    return {
     'status': 'fail',
     'error': 'Unable to extract relevant content from corpus,  Either there is less than 20 sentences in corpus for this topic, or structure of corpus is poor' 
    }

  #call a function here that make ppt and return its location
  #replace ppt_loc value from actual location of ppt
  ppt_loc = generate_ppt(content)

  return {
    'status': 'success',
    'ppt': ppt_loc
  }
